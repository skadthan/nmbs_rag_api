from langchain.document_loaders import CSVLoader
from langchain.text_splitter import CharacterTextSplitter
from langchain.indexes.vectorstore import VectorStoreIndexWrapper
from langchain.vectorstores import FAISS
from pathlib import Path

#from langchain.embeddings import BedrockEmbeddings #Warning message, this is updated to the the following import.
from langchain_aws import BedrockEmbeddings
from langchain_community.document_loaders import Docx2txtLoader
from langchain.document_loaders import PyPDFLoader, PyPDFDirectoryLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from pptx import Presentation
from openpyxl import load_workbook
from langchain.schema import Document
from app.utilities.llm_client import get_bedrock_embedding_model
from app.utilities.esclient import get_es_connection
from app.utilities.s3.download_file import sync_s3_bucket_to_local
import os
from app import config
from langchain_elasticsearch import ElasticsearchStore
from pathlib import Path
import boto3
from botocore.config import Config
from opensearchpy import OpenSearch, RequestsHttpConnection, AWSV4SignerAuth
from langchain_community.vectorstores import OpenSearchVectorSearch

def load_docx(file):
    loader = Docx2txtLoader(file)
    documents= loader.load()
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=2000,chunk_overlap=200,)
    split_data = text_splitter.split_documents(documents)
    #print(f"Number of documents={len(documents)}")
    return split_data

def load_pdf(file):
    loader = PyPDFDirectoryLoader(file)
    documents = loader.load()
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=2000,chunk_overlap=200,)
    split_data = text_splitter.split_documents(documents)
    #print(f"Number of documents={len(documents)}")
    return split_data

def load_csv(file):
    loader = CSVLoader(file) # --- > 219 docs with 400 chars, each row consists in a question column and an answer column
    documents = loader.load() #
    #print(f"Number of documents={len(documents)}")
    split_data = CharacterTextSplitter(chunk_size=2000, chunk_overlap=400, separator=",").split_documents(documents)
    return split_data

def load_xlsx(file):
    workbook = load_workbook(file)
    full_text = ""
    documents = []
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join([str(cell) if cell is not None else "" for cell in row])
            full_text += row_text + "\n"
    
    documents.append(Document(page_content=full_text, metadata={"source": str(file)}))
    # Split the text into chunks
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=2000,chunk_overlap=200,)
    split_data = text_splitter.split_documents(documents)
    #print(f"Processed {len(split_data)} document chunks.")
    return split_data

def load_pptx(file):
    presentation = Presentation(file)
    full_text = ""
    documents = []
    
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                full_text += shape.text + "\n"
    
    documents.append(Document(page_content=full_text, metadata={"source": str(file)}))
    # Split the text into chunks
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=2000,chunk_overlap=200,)
    split_data = text_splitter.split_documents(documents)
    #print(f"Processed {len(split_data)} document chunks.")
    return split_data
    

def process_s3_bucket(bucket_name: str, vector_index: str):
    es_vector_store=None
    aoss_vector_store=None
    folder_path = config.S3_LOCAL_DATA_FOLDER
    s3_prefix=""
    sync_s3_bucket_to_local(bucket_name,folder_path,s3_prefix)

    if not vector_index:
        #aoss_es_vector_store_name=config.DEFAULT_INDEX_NAME #Local Elastic Search Default Index Name
        aoss_es_vector_store_name=config.AOSS_VECTORSTORE_NAME
    else:
        aoss_es_vector_store_name=vector_index
    
    embedding_model = get_bedrock_embedding_model()
    es_client = get_es_connection()
    folder_path = Path(folder_path)
    try:
        os.makedirs(folder_path, exist_ok=True)
        #print(f"Directory created or already exists: {folder_path}")
    except Exception as e:
        print(f"An error occurred while creating the directory: {e}")


    for file_path in folder_path.iterdir():
        #print("Printing file path--> ", file_path)
        try:
            if file_path.is_file():
                #print(f"Processing file: {file_path}")
                file_extension = os.path.splitext(file_path)[1].lower()
                #print("printing file_extension --> ", file_extension)
                # Check the type based on extension
                
                if file_extension == '.docx':
                    file_type = 'Word Document'
                    docs=load_docx(file_path)
                elif file_extension == '.ppt' or file_extension == '.pptx':
                    file_type = 'PowerPoint Presentation'
                    docs=load_pptx(file_path)
                elif file_extension == '.xlsx':
                    file_type = 'Excel Spreadsheet'
                    docs=load_xlsx(file_path)
                elif file_extension == '.pdf':
                    file_type = 'PDF Document'
                    loader = PyPDFDirectoryLoader(file_path)
                    docs=load_pdf(file_path)
                elif file_extension == '.csv':
                    file_type = 'CSV File'
                    docs=load_csv(file_path)
                else:
                    file_type = 'Unknown Type'

                #print(f"{file_path}: {file_type}")

                # Update or create the ElasticSearch vectorstore
                #es_vector_store = store_opensearch_embeddings(docs,embedding_model,es_vector_store,es_client,aoss_es_vector_store_name)
                
                #Also store these embeddings in Amazon OpenSearch Serverless.
                aoss_vector_store = store_opensearch_embeddings(docs,embedding_model,aoss_vector_store,aoss_es_vector_store_name)

        except Exception as e:
            print(f"Error processing {file_path}: {e}")


def store_opensearch_embeddings(docs,embedding_model,aoss_vector_store,aoss_es_vector_store_name):
    #print("in the fuction --> store_opensearch_embeddings")
    vector_store_name = aoss_es_vector_store_name #config.AOSS_VECTORSTORE_NAME
    index_name = config.AOSS_VECTORSTORE_INDEX
   
    aoss_client = boto3.client('opensearchserverless',config=Config(region_name=config.AWS_REGION))
    collection = aoss_client.batch_get_collection(names=[vector_store_name])
    #print("Printing collection name", collection)
    host = collection['collectionDetails'][0]['id'] + '.' + config.AWS_REGION + '.aoss.amazonaws.com:443'
    #print("SURESH-OpenSource-Host", host)

    service = 'aoss'
    credentials = boto3.Session().get_credentials()
    auth = AWSV4SignerAuth(credentials, config.AWS_REGION, service)

    # Update or create the ElasticSearch vectorstore
    bulk_size = max(len(docs), 1000)
    if aoss_vector_store is None:
        aoss_vector_store = OpenSearchVectorSearch.from_documents(
        docs,
        embedding_model,
        opensearch_url=host,
        http_auth=auth,
        timeout = 100,
        use_ssl = True,
        verify_certs = True,
        connection_class = RequestsHttpConnection,
        index_name=index_name,
        engine="faiss",
        bulk_size = bulk_size
    )
    else:
        #print("In Else: Adding more documents...")
        aoss_vector_store.add_documents(documents=docs)
    return aoss_vector_store

def store_elasticsearch_embeddings(docs,embedding_model,es_vector_store,es_client,aoss_es_vector_store_name):
    if es_vector_store is None:
        es_vector_store = ElasticsearchStore(es_connection=es_client,index_name=aoss_es_vector_store_name,embedding=embedding_model)
        es_vector_store = es_vector_store.from_documents(documents=docs,embedding=embedding_model,index_name=aoss_es_vector_store_name,es_connection=es_client)
    else:
        es_vector_store.add_documents(documents=docs)
    return es_vector_store
